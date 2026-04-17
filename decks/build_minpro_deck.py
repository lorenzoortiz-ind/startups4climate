"""
S4C — Pitch deck para Ministerio de la Producción (MINPRO)
Genera un .pptx en formato 16:9 con la marca Startups4Climate.

Output: decks/S4C_MINPRO_pitch.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Brand palette ───────────────────────────────────────────────
BG_DARK     = RGBColor(0x07, 0x09, 0x0E)   # casi negro
BG_PRIMARY  = RGBColor(0x0B, 0x0E, 0x14)   # fondo base
BG_CARD     = RGBColor(0x13, 0x17, 0x1F)   # tarjetas
BG_ELEVATED = RGBColor(0x1E, 0x23, 0x30)   # capa elevada
ACCENT      = RGBColor(0xFF, 0x6B, 0x4A)   # coral S4C
ACCENT_2    = RGBColor(0xFF, 0x82, 0x66)   # coral hover
ACCENT_TEAL = RGBColor(0x0D, 0x94, 0x88)   # secundario (organizaciones)
ACCENT_AMB  = RGBColor(0xF5, 0x9E, 0x0B)   # secundario (superadmin)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # secundario (data)
ACCENT_GRN  = RGBColor(0x10, 0xB9, 0x81)   # secundario (success)
RED_BRAND   = RGBColor(0xE6, 0x39, 0x46)   # Redesign Lab
TEXT_HIGH   = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_MED    = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_LOW    = RGBColor(0x9C, 0xA3, 0xAF)
TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)
BORDER      = RGBColor(0x2A, 0x2F, 0x3D)

# ─── Slide geometry (16:9) ───────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)


# ────────────────────────────────────────────────────────────────
#  Helpers
# ────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color=BG_PRIMARY):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None, corner=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    if corner:
        # adjust corner radius (0 = no round, 0.5 = max)
        s.adjustments[0] = corner
    return s


def add_text_box(slide, x, y, w, h, text, *,
                 size=14, bold=False, color=TEXT_HIGH,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font='Inter', italic=False, line_spacing=1.2):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tx


def add_eyebrow(slide, x, y, text, color=ACCENT, size=10):
    add_text_box(slide, x, y, Inches(8), Inches(0.3),
                 text.upper(), size=size, bold=True, color=color, font='Inter')


def add_accent_bar(slide, x, y, w=Inches(0.55), h=Pt(3), color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def add_footer(slide, idx, total):
    add_text_box(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
                 'Startups4Climate · Pitch MINPRO 2026',
                 size=9, color=TEXT_MUTED)
    add_text_box(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
                 f'{idx} / {total}', size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_logo_block(slide, x, y, scale=1.0):
    """Visual S4C logo block: 'S4C' badge + Startups4Climate wordmark."""
    bs = Inches(0.42 * scale)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bs, bs)
    badge.fill.solid(); badge.fill.fore_color.rgb = BG_ELEVATED
    badge.line.color.rgb = ACCENT
    badge.line.width = Pt(0.75)
    badge.adjustments[0] = 0.25
    add_text_box(slide, x, y, bs, bs,
                 'S4C', size=int(11 * scale), bold=True, color=TEXT_HIGH,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font='Inter')
    # tiny accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + bs - Inches(0.07 * scale),
                                 y + Inches(0.05 * scale),
                                 Inches(0.075 * scale),
                                 Inches(0.075 * scale))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    # wordmark
    w_box = slide.shapes.add_textbox(x + bs + Inches(0.18 * scale),
                                     y - Inches(0.02 * scale),
                                     Inches(3.5 * scale), bs + Inches(0.04 * scale))
    tf = w_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    runs = [('Startups', TEXT_HIGH, False),
            ('4', ACCENT, False),
            ('Climate', TEXT_HIGH, False)]
    for txt, col, _ in runs:
        r = p.add_run(); r.text = txt
        r.font.name = 'Inter'; r.font.size = Pt(int(15 * scale))
        r.font.bold = True; r.font.color.rgb = col


def add_metric_card(slide, x, y, w, h, value, label, *,
                    accent=ACCENT, value_size=32, label_size=10):
    add_rect(slide, x, y, w, h, BG_CARD, line=BORDER, corner=0.06)
    add_accent_bar(slide, x + Inches(0.25), y + Inches(0.32),
                   w=Inches(0.32), h=Pt(2.5), color=accent)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.5),
                 w - Inches(0.5), Inches(0.7), value,
                 size=value_size, bold=True, color=TEXT_HIGH, font='Inter',
                 line_spacing=1.0)
    add_text_box(slide, x + Inches(0.25), y + h - Inches(0.55),
                 w - Inches(0.5), Inches(0.4), label,
                 size=label_size, color=TEXT_LOW, font='Inter', line_spacing=1.25)


def add_feature_card(slide, x, y, w, h, title, body, accent=ACCENT, kicker=None):
    add_rect(slide, x, y, w, h, BG_CARD, line=BORDER, corner=0.06)
    cy = y + Inches(0.35)
    if kicker:
        add_text_box(slide, x + Inches(0.35), cy, w - Inches(0.7), Inches(0.25),
                     kicker.upper(), size=9, bold=True, color=accent, font='Inter')
        cy += Inches(0.32)
    add_accent_bar(slide, x + Inches(0.35), cy + Inches(0.02),
                   w=Inches(0.4), h=Pt(2.5), color=accent)
    cy += Inches(0.18)
    add_text_box(slide, x + Inches(0.35), cy, w - Inches(0.7), Inches(0.5), title,
                 size=15, bold=True, color=TEXT_HIGH, font='Inter', line_spacing=1.15)
    add_text_box(slide, x + Inches(0.35), cy + Inches(0.55),
                 w - Inches(0.7), h - (cy - y) - Inches(0.85), body,
                 size=11, color=TEXT_LOW, font='Inter', line_spacing=1.4)


def add_bullet_list(slide, x, y, w, h, items, *, size=12, color=TEXT_MED, gap=0.18):
    cy = y
    for txt in items:
        # bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x, cy + Inches(0.1),
                                     Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text_box(slide, x + Inches(0.2), cy,
                     w - Inches(0.2), Inches(0.5),
                     txt, size=size, color=color, font='Inter', line_spacing=1.4)
        cy += Inches(gap + 0.18)


# ────────────────────────────────────────────────────────────────
#  Slide builders
# ────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slides = []

    # ───── 1. Cover ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s, BG_DARK)
    # Glow circle decoration
    glow1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(9.5), Inches(-1.5),
                               Inches(6), Inches(6))
    glow1.fill.solid(); glow1.fill.fore_color.rgb = RGBColor(0x2A, 0x12, 0x0A)
    glow1.line.fill.background()
    glow2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(-2), Inches(4.5),
                               Inches(5), Inches(5))
    glow2.fill.solid(); glow2.fill.fore_color.rgb = RGBColor(0x12, 0x14, 0x1A)
    glow2.line.fill.background()
    add_logo_block(s, Inches(0.7), Inches(0.55), scale=1.1)
    add_eyebrow(s, Inches(0.7), Inches(2.3),
                'Propuesta para el Ministerio de la Producción')
    add_text_box(s, Inches(0.7), Inches(2.65), Inches(11), Inches(2.0),
                 ['Infraestructura digital nacional',
                  'para programas de innovación y emprendimiento.'],
                 size=44, bold=True, color=TEXT_HIGH, font='Inter',
                 line_spacing=1.05)
    add_accent_bar(s, Inches(0.7), Inches(4.85),
                   w=Inches(0.7), h=Pt(3), color=ACCENT)
    add_text_box(s, Inches(0.7), Inches(5.0), Inches(10), Inches(1.0),
                 'Una plataforma todo-en-uno que estandariza la gestión de '
                 'cohortes, mide impacto y conecta el ecosistema peruano de '
                 'startups de impacto.',
                 size=15, color=TEXT_LOW, font='Inter', line_spacing=1.45)
    # bottom credits
    add_text_box(s, Inches(0.7), Inches(6.55), Inches(7), Inches(0.4),
                 'Desarrollado por Redesign Lab · Lima, Perú · Abril 2026',
                 size=11, color=TEXT_MUTED, font='Inter')
    add_text_box(s, Inches(8.6), Inches(6.55), Inches(4.2), Inches(0.4),
                 'startups4climate.vercel.app',
                 size=11, color=ACCENT, font='Inter', align=PP_ALIGN.RIGHT)

    # ───── 2. Contexto / Oportunidad ──────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '01 · Contexto')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'El Perú tiene talento emprendedor — pero infraestructura fragmentada.',
                 size=28, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(2.0), Inches(12), Inches(0.6),
                 'Tres realidades que limitan la escala del emprendimiento de impacto.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)
    cards = [
        ('Talento disperso',
         'Más de 200 incubadoras, aceleradoras y programas universitarios — sin '
         'estándares comunes de medición ni base de datos compartida.'),
        ('Programas sin trazabilidad',
         'Concursos como StartUp Perú, ProInnóvate y InnóvateMype mueven millones '
         'soles anualmente, pero el seguimiento post-financiamiento es manual y opaco.'),
        ('Founders con poca preparación',
         'Solo 1 de cada 5 startups que reciben financiamiento público sobrevive '
         'al año 3. La diferencia: fundamentos de negocio, no la idea.'),
    ]
    cw = Inches(3.95); cx = MARGIN
    for title, body in cards:
        add_feature_card(s, cx, Inches(2.8), cw, Inches(3.3),
                         title, body, accent=ACCENT)
        cx += cw + Inches(0.18)
    add_footer(s, 2, 19)

    # ───── 3. El problema (data sharper) ──────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '02 · El problema')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Cada programa reinventa la rueda. Y nadie ve el cuadro completo.',
                 size=28, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    # left list
    add_text_box(s, MARGIN, Inches(2.1), Inches(6), Inches(0.5),
                 'Lo que pasa hoy en cada programa:', size=13, bold=True,
                 color=TEXT_HIGH)
    pains = [
        'Excel + WhatsApp + correos como sistema de gestión.',
        'Diagnóstico inicial diferente en cada incubadora.',
        'Reportes manuales que tardan semanas en consolidarse.',
        'Founders sin acceso continuo a herramientas estructuradas.',
        'Imposible comparar resultados entre cohortes o regiones.',
        'Data perdida cuando termina el programa.',
    ]
    add_bullet_list(s, MARGIN, Inches(2.55), Inches(6), Inches(4),
                    pains, size=12, gap=0.22)
    # right metric card stack
    rx = Inches(7.3)
    add_metric_card(s, rx, Inches(2.1), Inches(5.4), Inches(1.2),
                    '~70%', 'de programas no reportan tracking de '
                    'KPIs post-graduación', accent=ACCENT, value_size=26)
    add_metric_card(s, rx, Inches(3.4), Inches(5.4), Inches(1.2),
                    '> 6 meses', 'tarda un programa en preparar el reporte '
                    'final consolidado', accent=ACCENT_AMB, value_size=24)
    add_metric_card(s, rx, Inches(4.7), Inches(5.4), Inches(1.2),
                    'S/. 0', 'visibilidad consolidada del impacto que financia '
                    'el Estado', accent=ACCENT_BLUE, value_size=26)
    add_metric_card(s, rx, Inches(6.0), Inches(5.4), Inches(0.95),
                    '1 de 5', 'startups públicas sobrevive al año 3',
                    accent=ACCENT_TEAL, value_size=22)
    add_footer(s, 3, 19)

    # ───── 4. La propuesta ────────────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '03 · Nuestra propuesta')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Startups4Climate: el sistema operativo del ecosistema.',
                 size=28, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(2.0), Inches(12), Inches(1.0),
                 'Una sola plataforma con tres capas integradas. Cada usuario ve lo '
                 'suyo. El Ministerio ve todo el sistema en tiempo real.',
                 size=14, color=TEXT_LOW, line_spacing=1.45)
    # 3 capas visualmente
    layers = [
        ('Capa 1', 'Founders', 'Herramientas + AI mentor.\n'
         'Diagnóstico, plan de negocios, finanzas, pitch — '
         'todo guiado y exportable.', ACCENT),
        ('Capa 2', 'Organizaciones', 'Incubadoras, aceleradoras, universidades.\n'
         'Cohortes, invitaciones, dashboard, reportes Excel y benchmarking.',
         ACCENT_TEAL),
        ('Capa 3', 'MINPRO / Ente rector', 'Vista nacional comparativa.\n'
         'Programas, regiones, verticales, presupuesto, hitos, impacto y ROI.',
         ACCENT_AMB),
    ]
    cw = Inches(3.95); cx = MARGIN; cy = Inches(3.6)
    for kicker, title, body, accent in layers:
        add_feature_card(s, cx, cy, cw, Inches(3.2), title, body,
                         accent=accent, kicker=kicker)
        cx += cw + Inches(0.18)
    add_footer(s, 4, 19)

    # ───── 5. Quiénes somos — Redesign Lab ───────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '04 · Quiénes somos')
    # Re.design Lab styled heading
    title_box = s.shapes.add_textbox(MARGIN, Inches(0.85), Inches(11), Inches(1.2))
    tf = title_box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0); tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.line_spacing = 1.05
    parts = [('Construido por ', TEXT_HIGH), ('Re', TEXT_HIGH),
             ('.', RED_BRAND), ('design ', TEXT_HIGH), ('Lab', RED_BRAND)]
    for t, c in parts:
        r = p.add_run(); r.text = t
        r.font.name = 'Inter'; r.font.size = Pt(36)
        r.font.bold = True; r.font.color.rgb = c
    add_text_box(s, MARGIN, Inches(2.15), Inches(12), Inches(1.0),
                 'Estudio peruano de innovación con foco en infraestructura digital '
                 'para emprendimiento de impacto en Latinoamérica.',
                 size=14, color=TEXT_LOW, line_spacing=1.45)

    # 4 quick stats
    stats = [
        ('5+ años', 'diseñando programas de innovación'),
        ('3', 'universidades clientes activas'),
        ('17+', 'founders con cuenta activa en piloto'),
        ('LATAM', 'foco regional · sede en Lima'),
    ]
    cw = Inches(3.0); cx = MARGIN
    for v, l in stats:
        add_metric_card(s, cx, Inches(3.4), cw, Inches(1.4),
                        v, l, accent=RED_BRAND,
                        value_size=22, label_size=10)
        cx += cw + Inches(0.13)

    # mission card
    add_rect(s, MARGIN, Inches(5.1), Inches(12.13), Inches(1.7),
             BG_ELEVATED, line=BORDER, corner=0.04)
    add_text_box(s, MARGIN + Inches(0.4), Inches(5.3),
                 Inches(11.3), Inches(0.4),
                 'NUESTRA MISIÓN', size=10, bold=True, color=ACCENT)
    add_text_box(s, MARGIN + Inches(0.4), Inches(5.55),
                 Inches(11.3), Inches(1.3),
                 '"Democratizar el desarrollo de startups de impacto en Latinoamérica '
                 'construyendo la infraestructura operativa que founders, organizaciones '
                 'y entes públicos necesitan para dejar de improvisar."',
                 size=14, italic=True, color=TEXT_HIGH, line_spacing=1.4)
    add_footer(s, 5, 19)

    # ───── 6. Tracción actual ────────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '05 · Tracción')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Ya operativos en 3 universidades — antes de marketing externo.',
                 size=28, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(2.0), Inches(12), Inches(0.6),
                 'Llevamos meses validando con instituciones reales. Esto es lo que vemos.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)

    # logo orgs row
    orgs = [
        ('UNAMAD', 'Univ. Nacional Amazónica de Madre de Dios', ACCENT_GRN),
        ('Wiener', 'Universidad Privada Norbert Wiener', ACCENT_BLUE),
        ('BioInnova', 'Programa de bioemprendimiento', ACCENT),
    ]
    cw = Inches(3.95); cx = MARGIN
    for n, sub, acc in orgs:
        add_rect(s, cx, Inches(2.85), cw, Inches(1.5),
                 BG_CARD, line=BORDER, corner=0.05)
        add_accent_bar(s, cx + Inches(0.35), Inches(3.05),
                       w=Inches(0.4), h=Pt(2.5), color=acc)
        add_text_box(s, cx + Inches(0.35), Inches(3.2),
                     cw - Inches(0.7), Inches(0.5), n,
                     size=20, bold=True, color=TEXT_HIGH)
        add_text_box(s, cx + Inches(0.35), Inches(3.65),
                     cw - Inches(0.7), Inches(0.55), sub,
                     size=11, color=TEXT_LOW, line_spacing=1.4)
        cx += cw + Inches(0.18)

    # bottom metrics row
    bx = MARGIN
    bw = Inches(3.0)
    bottom = [
        ('30+', 'Herramientas funcionales', ACCENT),
        ('200', 'Founders activos en pipeline', ACCENT_TEAL),
        ('30', 'Universidades target Q3 2026', ACCENT_AMB),
        ('99.9%', 'Uptime últimos 30 días', ACCENT_GRN),
    ]
    for v, l, c in bottom:
        add_metric_card(s, bx, Inches(4.7), bw, Inches(1.5),
                        v, l, accent=c, value_size=26)
        bx += bw + Inches(0.13)

    # footer note
    add_text_box(s, MARGIN, Inches(6.45), Inches(12), Inches(0.4),
                 'Producción estable en startups4climate.vercel.app · '
                 'Plan de escala: 200 founders + 30 universidades en 30 días.',
                 size=11, italic=True, color=TEXT_MUTED)
    add_footer(s, 6, 19)

    # ───── 7. Capa Founder ───────────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '06 · Capa 1 — Founders')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Una caja de herramientas profesional que reemplaza 30+ servicios.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(1.95), Inches(12), Inches(0.7),
                 'Acceso gratuito para founders. Cada herramienta es accionable, '
                 'guarda progreso y exporta reportes en PDF.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)

    feats = [
        ('Diagnóstico', 'Startup Readiness Score con 60+ preguntas y plan de acción.'),
        ('Plan de negocios', 'Lean Canvas, Business Model, Value Proposition Canvas.'),
        ('Finanzas', 'Unit economics, proyecciones, cap table, runway, fundraising.'),
        ('Mercado', 'TAM/SAM/SOM, JTBD, segmentación, MVP roadmap.'),
        ('Pitch', 'Pitch deck guiado, narrativa, demo day prep.'),
        ('Mentor AI', 'Gemini 2.5 Flash con contexto del founder y vertical.'),
        ('Passport', 'Identidad pública compartible con investors y aceleradoras.'),
        ('Reportes', 'Exportables PDF con marca de la organización.'),
    ]
    cols = 4
    cw = Inches(3.0); ch = Inches(1.6); cgap = Inches(0.13)
    for i, (t, body) in enumerate(feats):
        row = i // cols; col = i % cols
        x = MARGIN + col * (cw + cgap)
        y = Inches(2.95) + row * (ch + cgap)
        add_rect(s, x, y, cw, ch, BG_CARD, line=BORDER, corner=0.06)
        add_text_box(s, x + Inches(0.25), y + Inches(0.2),
                     cw - Inches(0.5), Inches(0.4), t,
                     size=13, bold=True, color=TEXT_HIGH)
        add_text_box(s, x + Inches(0.25), y + Inches(0.6),
                     cw - Inches(0.5), ch - Inches(0.7), body,
                     size=10, color=TEXT_LOW, line_spacing=1.4)
    add_footer(s, 7, 19)

    # ───── 8. Fundamento metodológico ─────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55),
                '07 · Fundamento metodológico')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12.2), Inches(1.0),
                 'Cada herramienta está construida sobre marcos académicos validados.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(1.95), Inches(12.2), Inches(0.7),
                 'No reinventamos la rueda: sintetizamos lo mejor de MIT, Stanford, '
                 'Strategyzer, Y Combinator, GIIN y otros — y lo hacemos accionable '
                 'en español, gratis para founders LATAM.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)

    # 8 methodology cards — color-coded by source
    methods = [
        ('MIT', ACCENT,
         'Disciplined Entrepreneurship — 24 Steps',
         'Bill Aulet · MIT Martin Trust Center',
         'Beachhead Market, End User Profile, Persona, '
         'Full Lifecycle Use Case, QVP, DMU, Core, MVBP.'),
        ('STANFORD', ACCENT_BLUE,
         'Customer Development & Lean Startup',
         'Steve Blank · Eric Ries',
         'Key Assumptions, MVP Definition, Traction '
         'Validation, First 10 Customers.'),
        ('STRATEGYZER', ACCENT_TEAL,
         'Business Model & Value Prop Canvas',
         'Alexander Osterwalder · Yves Pigneur',
         'Business Model Design, Quantified Value Prop, '
         'Pricing Framework.'),
        ('LEANSTACK', ACCENT_TEAL,
         'Lean Canvas',
         'Ash Maurya',
         'Lean Canvas — versión 1-página optimizada para '
         'startups en stage temprano.'),
        ('STANFORD d.school', ACCENT_BLUE,
         'Design Thinking & Purpose-Driven Founding',
         'Stanford d.school · IDEO',
         'Passion & Purpose, framing del problema, '
         'discovery del founder.'),
        ('Y COMBINATOR · SEQUOIA', ACCENT_AMB,
         'Pitch & Fundraising Frameworks',
         'YC Demo Day · Sequoia Memo · a16z',
         'Pitch Deck guiado, narrativa de inversión, '
         'cap table y dilución.'),
        ('INTEL · GOOGLE', ACCENT_AMB,
         'OKRs — Objectives & Key Results',
         'Andy Grove · John Doerr (Measure What Matters)',
         'OKR Tracker para founders y cohortes.'),
        ('GIIN · IRIS+ · ODS', ACCENT_GRN,
         'Estándares de impacto',
         'Global Impact Investing Network · IRIS+ · '
         'SDG Compass de Naciones Unidas',
         'Impact Metrics, Theory of Change, alineación a '
         'los 17 ODS.'),
    ]

    cols = 4
    cw = Inches(3.0); ch = Inches(1.85); cgap = Inches(0.13)
    base_y = Inches(2.95)
    for i, (kicker, kc, title, author, body) in enumerate(methods):
        row = i // cols; col = i % cols
        x = MARGIN + col * (cw + cgap)
        y = base_y + row * (ch + cgap)
        add_rect(s, x, y, cw, ch, BG_CARD, line=BORDER, corner=0.06)
        # color stripe (left edge)
        add_accent_bar(s, x + Inches(0.22), y + Inches(0.18),
                       w=Inches(0.42), h=Pt(2.5), color=kc)
        # kicker (institution)
        add_text_box(s, x + Inches(0.22), y + Inches(0.22),
                     cw - Inches(0.44), Inches(0.25),
                     kicker, size=8, bold=True, color=kc, font='Inter')
        # title (methodology)
        add_text_box(s, x + Inches(0.22), y + Inches(0.5),
                     cw - Inches(0.44), Inches(0.5),
                     title, size=12, bold=True, color=TEXT_HIGH,
                     font='Inter', line_spacing=1.15)
        # author
        add_text_box(s, x + Inches(0.22), y + Inches(0.95),
                     cw - Inches(0.44), Inches(0.32),
                     author, size=9, italic=True, color=TEXT_LOW,
                     font='Inter', line_spacing=1.25)
        # description
        add_text_box(s, x + Inches(0.22), y + Inches(1.22),
                     cw - Inches(0.44), ch - Inches(1.32),
                     body, size=9, color=TEXT_MED,
                     font='Inter', line_spacing=1.35)

    # bottom note
    add_text_box(s, MARGIN, Inches(6.55), Inches(12.2), Inches(0.4),
                 'Cada herramienta cita su fuente metodológica dentro de la '
                 'plataforma · Las versiones en español son adaptación propia '
                 'de Redesign Lab.',
                 size=10, italic=True, color=TEXT_MUTED, line_spacing=1.4)
    add_footer(s, 8, 19)

    # ───── 9. Capa Organización ──────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '07 · Capa 2 — Organizaciones',
                color=ACCENT_TEAL)
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Gestión integral de cohortes — sin más Excel ni WhatsApp.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(1.95), Inches(12), Inches(0.7),
                 'Para incubadoras, aceleradoras, universidades y operadores de '
                 'programas públicos. Plan Regional gratuito; Profesional desde S/ 1,100/año.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)

    org_feats = [
        ('Cohortes', 'Crear cohortes con fechas, milestones y descripción.'),
        ('Invitaciones', 'Founders por email, registro automático, vinculación.'),
        ('Dashboard', 'Progreso de cada startup en tiempo real.'),
        ('Benchmarking', 'Comparar cohortes vs. promedio plataforma.'),
        ('Reportes', 'Excel descargable listo para stakeholders.'),
        ('Configuración', 'Branding, logo, hitos personalizados.'),
    ]
    cols = 3
    cw = Inches(4.04); ch = Inches(1.55); cgap = Inches(0.14)
    for i, (t, body) in enumerate(org_feats):
        row = i // cols; col = i % cols
        x = MARGIN + col * (cw + cgap)
        y = Inches(2.95) + row * (ch + cgap)
        add_rect(s, x, y, cw, ch, BG_CARD, line=BORDER, corner=0.06)
        add_accent_bar(s, x + Inches(0.3), y + Inches(0.25),
                       w=Inches(0.32), h=Pt(2.5), color=ACCENT_TEAL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.4),
                     cw - Inches(0.6), Inches(0.4), t,
                     size=14, bold=True, color=TEXT_HIGH)
        add_text_box(s, x + Inches(0.3), y + Inches(0.8),
                     cw - Inches(0.6), ch - Inches(0.9), body,
                     size=11, color=TEXT_LOW, line_spacing=1.4)
    add_footer(s, 9, 19)

    # ───── 10. Capa Superadmin / MINPRO ────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '08 · Capa 3 — Vista nacional (MINPRO)',
                color=ACCENT_AMB)
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Lo que el MINPRO necesita: un dashboard para gobernar el sistema.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(1.95), Inches(12), Inches(0.7),
                 'Acceso exclusivo del rol Superadmin. Comparar programas, regiones, '
                 'verticales, presupuesto, ROI e impacto — en tiempo real.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)
    panels = [
        ('Comparar programas', 'StartUp Perú, ProInnóvate, InnóvateMype y operadores '
         'regionales — métricas estandarizadas.'),
        ('Por región', '24 departamentos con cobertura, número de founders, score '
         'promedio y ejecución presupuestal.'),
        ('Por vertical', 'Cleantech, agritech, healthtech, edtech, fintech — '
         'distribución por etapa y madurez.'),
        ('Presupuesto y ROI', 'Inversión por founder, costo por startup graduada, '
         'comparable entre programas.'),
        ('Hitos y avance', 'Milestones por programa, retrasos, alertas tempranas.'),
        ('Impacto agregado', 'Empleos generados, capital movilizado, tCO₂e, '
         'beneficiarios indirectos.'),
    ]
    cols = 3
    cw = Inches(4.04); ch = Inches(1.7); cgap = Inches(0.14)
    for i, (t, body) in enumerate(panels):
        row = i // cols; col = i % cols
        x = MARGIN + col * (cw + cgap)
        y = Inches(2.95) + row * (ch + cgap)
        add_rect(s, x, y, cw, ch, BG_CARD, line=BORDER, corner=0.06)
        add_accent_bar(s, x + Inches(0.3), y + Inches(0.25),
                       w=Inches(0.32), h=Pt(2.5), color=ACCENT_AMB)
        add_text_box(s, x + Inches(0.3), y + Inches(0.4),
                     cw - Inches(0.6), Inches(0.4), t,
                     size=14, bold=True, color=TEXT_HIGH)
        add_text_box(s, x + Inches(0.3), y + Inches(0.8),
                     cw - Inches(0.6), ch - Inches(0.9), body,
                     size=11, color=TEXT_LOW, line_spacing=1.4)
    add_footer(s, 10, 19)

    # ───── 11. Casos de uso para MINPRO ───────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '09 · Cómo conecta con la cartera del MINPRO')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Programas existentes que pueden potenciarse desde el día 1.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    progs = [
        ('StartUp Perú',
         'Estandarizar diagnóstico, dashboard de cohorts y reportes para los '
         'operadores. Trazabilidad post-financiamiento.'),
        ('ProInnóvate',
         'Seguimiento de proyectos de innovación empresarial con KPIs unificados '
         'y comparables entre convocatorias.'),
        ('InnóvateMype',
         'Capacitación + herramientas estructuradas para Mypes en proceso de '
         'formalización y crecimiento.'),
        ('Programas regionales',
         'Cobertura para incubadoras de UNI, UNT, UNSA, UNALM, UNAS, UNMSM y '
         'gobiernos regionales — plan Regional gratuito.'),
    ]
    cw = Inches(6.06); ch = Inches(2.1); cgap = Inches(0.14)
    for i, (t, body) in enumerate(progs):
        row = i // 2; col = i % 2
        x = MARGIN + col * (cw + cgap)
        y = Inches(2.85) + row * (ch + cgap)
        add_feature_card(s, x, y, cw, ch, t, body, accent=ACCENT)
    add_footer(s, 11, 19)

    # ───── 12. Stack tecnológico ──────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '10 · Stack y arquitectura')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Construido con tecnología seria — pensado para el sector público.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    techs = [
        ('Next.js 15', 'Framework moderno, despliegue serverless en Vercel.'),
        ('Supabase', 'Postgres gestionado con RLS por organización y rol.'),
        ('Gemini 2.5 Flash', 'Motor AI para mentor contextualizado.'),
        ('Resend', 'Email transaccional para invitaciones y reportes.'),
        ('TypeScript', 'Tipado estricto, menor superficie de bugs.'),
        ('Vercel + GitHub', 'CI/CD, preview por branch, rollbacks 1-click.'),
    ]
    cols = 3
    cw = Inches(4.04); ch = Inches(1.4); cgap = Inches(0.14)
    for i, (t, body) in enumerate(techs):
        row = i // cols; col = i % cols
        x = MARGIN + col * (cw + cgap)
        y = Inches(2.85) + row * (ch + cgap)
        add_rect(s, x, y, cw, ch, BG_CARD, line=BORDER, corner=0.06)
        add_text_box(s, x + Inches(0.3), y + Inches(0.25),
                     cw - Inches(0.6), Inches(0.4), t,
                     size=13, bold=True, color=ACCENT)
        add_text_box(s, x + Inches(0.3), y + Inches(0.65),
                     cw - Inches(0.6), ch - Inches(0.7), body,
                     size=11, color=TEXT_LOW, line_spacing=1.4)

    # security stripe
    add_rect(s, MARGIN, Inches(5.85), Inches(12.13), Inches(1.0),
             BG_ELEVATED, line=BORDER, corner=0.04)
    add_text_box(s, MARGIN + Inches(0.35), Inches(6.0),
                 Inches(11.5), Inches(0.4),
                 'SEGURIDAD Y PRIVACIDAD',
                 size=10, bold=True, color=ACCENT_GRN)
    add_text_box(s, MARGIN + Inches(0.35), Inches(6.3),
                 Inches(11.5), Inches(0.6),
                 'Row Level Security en cada tabla · auth multi-rol · auditoría '
                 'de eventos · datos hosteados en infraestructura cloud certificada · '
                 'cumplimiento Ley 29733 (protección de datos personales).',
                 size=11, color=TEXT_MED, line_spacing=1.4)
    add_footer(s, 12, 19)

    # ───── 13. Tipos de usuario ───────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '11 · Tipos de usuario')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Cuatro perfiles, una sola plataforma.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)

    users = [
        ('Founder', 'Acceso gratuito. Usa herramientas, recibe mentoría AI, '
         'genera su Passport y se postula a oportunidades.', ACCENT,
         '17+ activos hoy'),
        ('Admin Org.', 'Líder de incubadora/programa. Crea cohortes, invita '
         'founders, descarga reportes, configura branding.', ACCENT_TEAL,
         '3 organizaciones piloto'),
        ('Superadmin', 'Rol exclusivo MINPRO o ente rector. Vista nacional '
         'comparativa de todos los programas.', ACCENT_AMB,
         '1 cuenta MINPRO propuesta'),
        ('Mentor (Q3)', 'Próximamente. Profesionales validados que pueden '
         'asignarse 1:1 a startups específicas.', ACCENT_BLUE,
         'Roadmap fase 2'),
    ]
    cw = Inches(3.0); ch = Inches(2.95)
    cx = MARGIN
    for t, body, acc, badge in users:
        add_rect(s, cx, Inches(2.85), cw, ch, BG_CARD, line=BORDER, corner=0.06)
        # Badge top
        add_rect(s, cx, Inches(2.85), cw, Inches(0.45), acc, corner=0.18)
        add_text_box(s, cx, Inches(2.93), cw, Inches(0.3), t.upper(),
                     size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.TOP)
        add_text_box(s, cx + Inches(0.25), Inches(3.5),
                     cw - Inches(0.5), Inches(1.7), body,
                     size=11, color=TEXT_MED, line_spacing=1.4)
        add_text_box(s, cx + Inches(0.25), Inches(5.3),
                     cw - Inches(0.5), Inches(0.4), 'ESTADO ACTUAL',
                     size=9, bold=True, color=TEXT_MUTED)
        add_text_box(s, cx + Inches(0.25), Inches(5.55),
                     cw - Inches(0.5), Inches(0.45), badge,
                     size=12, bold=True, color=acc)
        cx += cw + Inches(0.13)
    add_footer(s, 13, 19)

    # ───── 14. Impacto esperado ──────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '12 · Impacto en el ecosistema')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Lo que cambia cuando el ecosistema corre sobre infraestructura común.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    impacts = [
        ('Para el MINPRO', [
            'Visibilidad nacional en tiempo real de programas y operadores.',
            'Reportes de ejecución y resultados sin esperar 6 meses.',
            'Decisiones presupuestales basadas en data comparable.',
            'Detección temprana de programas sub-performantes.',
        ]),
        ('Para las organizaciones operadoras', [
            'Reducción 80% del tiempo administrativo por cohorte.',
            'Métricas estandarizadas — fácil comparar y mejorar.',
            'Reportes automatizados a stakeholders.',
            'Branding propio + acceso a comunidad nacional.',
        ]),
        ('Para los founders', [
            'Herramientas de calidad MIT/Stanford — gratis.',
            'Mentor AI 24/7 contextualizado a su vertical.',
            'Identidad pública (Passport) compartible con investors.',
            'Trazabilidad de su evolución a lo largo de programas.',
        ]),
    ]
    cw = Inches(4.04); ch = Inches(3.6); cgap = Inches(0.14); cx = MARGIN
    for t, items in impacts:
        add_rect(s, cx, Inches(2.85), cw, ch, BG_CARD, line=BORDER, corner=0.05)
        add_text_box(s, cx + Inches(0.3), Inches(3.05),
                     cw - Inches(0.6), Inches(0.5), t,
                     size=15, bold=True, color=ACCENT)
        cy = Inches(3.55)
        for it in items:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx + Inches(0.3), cy + Inches(0.08),
                                     Inches(0.07), Inches(0.07))
            dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
            dot.line.fill.background()
            add_text_box(s, cx + Inches(0.5), cy,
                         cw - Inches(0.8), Inches(0.6), it,
                         size=11, color=TEXT_MED, line_spacing=1.35)
            cy += Inches(0.62)
        cx += cw + cgap
    add_footer(s, 14, 19)

    # ───── 15. Propuesta de piloto ───────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '13 · Propuesta de piloto')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Piloto controlado de 6 meses con 3 programas y la cuenta MINPRO.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)

    # Left side: scope
    add_text_box(s, MARGIN, Inches(2.1), Inches(6), Inches(0.5),
                 'Alcance del piloto', size=14, bold=True, color=ACCENT)
    scope = [
        '3 organizaciones operadoras (1 limeña + 2 regionales).',
        '1 cuenta Superadmin para el equipo MINPRO designado.',
        '~150 founders vinculados a las cohortes.',
        'Acceso completo a 30+ herramientas y mentor AI.',
        'Onboarding presencial y soporte dedicado.',
        'Sesiones de levantamiento y retroalimentación quincenales.',
    ]
    add_bullet_list(s, MARGIN, Inches(2.55), Inches(6), Inches(4),
                    scope, size=12, gap=0.18)

    # Right: Timeline visual
    add_text_box(s, Inches(7.3), Inches(2.1), Inches(5.5), Inches(0.5),
                 'Timeline (6 meses)', size=14, bold=True, color=ACCENT)
    phases = [
        ('Mes 1', 'Setup', 'Onboarding orgs + cuenta MINPRO + import data', ACCENT),
        ('Mes 2-3', 'Operación', 'Cohortes activas, founders trabajando', ACCENT_TEAL),
        ('Mes 4-5', 'Optimización', 'Ajustes según feedback, primer reporte trimestral', ACCENT_AMB),
        ('Mes 6', 'Evaluación', 'Reporte final + decisión de escalado nacional', ACCENT_GRN),
    ]
    py = Inches(2.6)
    for m, ph, desc, c in phases:
        # Time chip
        add_rect(s, Inches(7.3), py, Inches(1.0), Inches(0.4), c, corner=0.25)
        add_text_box(s, Inches(7.3), py + Inches(0.04),
                     Inches(1.0), Inches(0.35), m,
                     size=10, bold=True, color=BG_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # Phase title + body
        add_text_box(s, Inches(8.45), py - Inches(0.05),
                     Inches(4.5), Inches(0.4), ph,
                     size=13, bold=True, color=TEXT_HIGH)
        add_text_box(s, Inches(8.45), py + Inches(0.32),
                     Inches(4.5), Inches(0.6), desc,
                     size=10, color=TEXT_LOW, line_spacing=1.4)
        py += Inches(0.95)
    add_footer(s, 15, 19)

    # ───── 16. KPIs del piloto ───────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '14 · KPIs y criterios de éxito')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Cómo medimos si el piloto funcionó.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    kpis = [
        ('80%', 'Founders que completan diagnóstico inicial', ACCENT),
        ('60%', 'Founders que terminan ≥3 herramientas', ACCENT_TEAL),
        ('100%', 'Reportes mensuales generados sin trabajo manual', ACCENT_AMB),
        ('NPS 50+', 'Satisfacción de los líderes de programa', ACCENT_BLUE),
        ('< 1 semana', 'Tiempo para preparar reporte ejecutivo', ACCENT_GRN),
        ('3', 'Programas en producción al cierre del piloto', ACCENT),
    ]
    cols = 3; cw = Inches(4.04); ch = Inches(1.6); cgap = Inches(0.14)
    for i, (v, l, c) in enumerate(kpis):
        row = i // cols; col = i % cols
        x = MARGIN + col * (cw + cgap)
        y = Inches(2.85) + row * (ch + cgap)
        add_metric_card(s, x, y, cw, ch, v, l, accent=c, value_size=28)

    # closing line
    add_rect(s, MARGIN, Inches(6.2), Inches(12.13), Inches(0.7),
             BG_ELEVATED, line=BORDER, corner=0.06)
    add_text_box(s, MARGIN + Inches(0.3), Inches(6.32),
                 Inches(11.5), Inches(0.5),
                 'Resultado esperado: data dura para decidir si escalar a nivel nacional en 2027.',
                 size=12, italic=True, bold=True, color=TEXT_HIGH, line_spacing=1.3)
    add_footer(s, 16, 19)

    # ───── 17. Modelo económico ──────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '15 · Modelo económico')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'Estructura de costos del piloto y proyección nacional.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    add_text_box(s, MARGIN, Inches(1.95), Inches(12), Inches(0.7),
                 'Modelo SaaS por organización + licencia institucional para el ente rector.',
                 size=13, color=TEXT_LOW, line_spacing=1.4)

    # Pricing cards
    plans = [
        ('Plan Regional', 'Gratuito',
         ['Hasta 20 startups', 'Acceso a herramientas',
          'Dashboard básico', 'Reportes simples',
          'Para incubadoras fuera de Lima'], ACCENT_TEAL, False),
        ('Plan Profesional', 'S/ 1,100 / año',
         ['Hasta 30 startups', 'Reportes avanzados Excel',
          'Benchmarking de cohortes', 'Soporte prioritario',
          'Branding personalizado'], ACCENT, True),
        ('Licencia MINPRO', 'A convenir',
         ['Cuenta Superadmin nacional', 'Vista de todos los programas',
          'API + integraciones', 'SLA dedicado',
          'Acompañamiento estratégico'], ACCENT_AMB, False),
    ]
    cw = Inches(4.04); cgap = Inches(0.14); cx = MARGIN
    for name, price, items, acc, hl in plans:
        ch = Inches(4.0)
        if hl:
            add_rect(s, cx, Inches(2.85), cw, ch, BG_ELEVATED,
                     line=acc, corner=0.06)
        else:
            add_rect(s, cx, Inches(2.85), cw, ch, BG_CARD,
                     line=BORDER, corner=0.06)
        add_text_box(s, cx + Inches(0.3), Inches(3.05),
                     cw - Inches(0.6), Inches(0.4), name,
                     size=14, bold=True, color=TEXT_HIGH)
        add_text_box(s, cx + Inches(0.3), Inches(3.45),
                     cw - Inches(0.6), Inches(0.6), price,
                     size=22, bold=True, color=acc)
        add_rect(s, cx + Inches(0.3), Inches(4.15), cw - Inches(0.6),
                 Emu(9525), BORDER)
        cy = Inches(4.3)
        for it in items:
            chk = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx + Inches(0.3), cy + Inches(0.08),
                                     Inches(0.08), Inches(0.08))
            chk.fill.solid(); chk.fill.fore_color.rgb = acc
            chk.line.fill.background()
            add_text_box(s, cx + Inches(0.5), cy,
                         cw - Inches(0.8), Inches(0.4), it,
                         size=11, color=TEXT_MED, line_spacing=1.35)
            cy += Inches(0.4)
        cx += cw + cgap
    add_footer(s, 17, 19)

    # ───── 18. Roadmap nacional ──────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s)
    add_eyebrow(s, MARGIN, Inches(0.55), '16 · Roadmap nacional post-piloto')
    add_text_box(s, MARGIN, Inches(0.85), Inches(12), Inches(1.0),
                 'De 3 organizaciones a cobertura nacional en 18 meses.',
                 size=26, bold=True, color=TEXT_HIGH, line_spacing=1.1)
    phases2 = [
        ('Q3 2026', 'Piloto cerrado',
         '3 orgs + 150 founders. Validación con el MINPRO. Reporte final.', ACCENT_GRN),
        ('Q4 2026', 'Onboarding macroregión',
         'Norte, Sur, Centro y Selva. 15 universidades + 5 incubadoras.', ACCENT_TEAL),
        ('Q1 2027', 'Integración cartera MINPRO',
         'StartUp Perú, ProInnóvate, InnóvateMype con dashboards comunes.', ACCENT_AMB),
        ('Q2 2027', 'Cobertura nacional',
         '24 departamentos · 30+ universidades · ente rector con visibilidad total.', ACCENT),
    ]
    py = Inches(2.85); pheight = Inches(1.0)
    for tag, t, body, c in phases2:
        # Time chip
        add_rect(s, MARGIN, py, Inches(1.5), Inches(0.4), c, corner=0.25)
        add_text_box(s, MARGIN, py + Inches(0.04),
                     Inches(1.5), Inches(0.35), tag,
                     size=11, bold=True, color=BG_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # Title + body
        add_text_box(s, Inches(2.35), py - Inches(0.05),
                     Inches(10.5), Inches(0.4), t,
                     size=15, bold=True, color=TEXT_HIGH)
        add_text_box(s, Inches(2.35), py + Inches(0.4),
                     Inches(10.5), Inches(0.6), body,
                     size=11, color=TEXT_LOW, line_spacing=1.4)
        py += pheight
    # Final note
    add_rect(s, MARGIN, Inches(6.4), Inches(12.13), Inches(0.6),
             BG_ELEVATED, line=BORDER, corner=0.06)
    add_text_box(s, MARGIN + Inches(0.3), Inches(6.5),
                 Inches(11.5), Inches(0.45),
                 'Visión 2030: S4C como infraestructura permanente del ecosistema '
                 'peruano de emprendimiento, replicable en LATAM.',
                 size=12, italic=True, bold=True, color=TEXT_HIGH)
    add_footer(s, 18, 19)

    # ───── 19. Cierre / Ask ──────────────────────────────────────
    s = prs.slides.add_slide(blank); slides.append(s)
    set_slide_bg(s, BG_DARK)
    # Decorations
    glow = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(8), Inches(-2),
                              Inches(7), Inches(7))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x2A, 0x12, 0x0A)
    glow.line.fill.background()

    add_logo_block(s, Inches(0.7), Inches(0.55), scale=0.95)

    add_eyebrow(s, MARGIN, Inches(2.0), 'El pedido')
    add_text_box(s, MARGIN, Inches(2.35), Inches(11.5), Inches(2.0),
                 ['Aprobar el piloto de 6 meses para validar S4C',
                  'como infraestructura del MINPRO.'],
                 size=36, bold=True, color=TEXT_HIGH, line_spacing=1.1)

    # Three asks
    asks = [
        ('1', 'Carta de respaldo institucional',
         'Para coordinar con 3 orgs operadoras del piloto.'),
        ('2', 'Designación de un equipo MINPRO',
         '2-3 personas como contraparte técnica y usuarios Superadmin.'),
        ('3', 'Definición de KPIs conjuntos',
         'Para evaluar el piloto y la decisión de escalado en Q1 2027.'),
    ]
    cw = Inches(4.04); cgap = Inches(0.14); cx = MARGIN
    for n, t, body in asks:
        add_rect(s, cx, Inches(4.7), cw, Inches(1.7),
                 BG_CARD, line=BORDER, corner=0.06)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx + Inches(0.3), Inches(4.9),
                                  Inches(0.42), Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        add_text_box(s, cx + Inches(0.3), Inches(4.92),
                     Inches(0.42), Inches(0.4), n,
                     size=14, bold=True, color=BG_DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        add_text_box(s, cx + Inches(0.85), Inches(4.92),
                     cw - Inches(1.2), Inches(0.4), t,
                     size=13, bold=True, color=TEXT_HIGH)
        add_text_box(s, cx + Inches(0.3), Inches(5.5),
                     cw - Inches(0.6), Inches(1.0), body,
                     size=11, color=TEXT_LOW, line_spacing=1.4)
        cx += cw + cgap

    # Contact stripe
    add_rect(s, MARGIN, Inches(6.65), Inches(12.13), Inches(0.55),
             BG_ELEVATED, line=BORDER, corner=0.06)
    add_text_box(s, MARGIN + Inches(0.3), Inches(6.75),
                 Inches(11.5), Inches(0.4),
                 'Lorenzo Ortiz · Co-fundador, Redesign Lab  ·  '
                 'lorenzo.ortiz@redesignlab.org  ·  WhatsApp +51 989 338 401  ·  '
                 'startups4climate.vercel.app',
                 size=10, color=TEXT_MED, align=PP_ALIGN.CENTER)
    return prs


if __name__ == '__main__':
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, 'S4C_MINPRO_pitch.pptx')
    prs = build()
    prs.save(out_path)
    print(f'OK · {out_path}')
    print(f'Slides: {len(prs.slides)}')
    size_kb = os.path.getsize(out_path) / 1024
    print(f'Size:   {size_kb:.1f} KB')
